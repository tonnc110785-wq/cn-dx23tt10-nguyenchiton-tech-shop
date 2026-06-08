#!/usr/bin/env python3
"""TechShop – Professional PowerPoint Generator (15 slides, dark-amber theme)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0F, 0x17, 0x2A)
BG2     = RGBColor(0x1E, 0x29, 0x3B)
BG3     = RGBColor(0x27, 0x37, 0x4F)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_D = RGBColor(0xD9, 0x77, 0x06)
AMBER_L = RGBColor(0xFC, 0xD3, 0x4D)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFF_W   = RGBColor(0xE2, 0xE8, 0xF0)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)
GREEN   = RGBColor(0x05, 0x96, 0x69)
GREEN_L = RGBColor(0x6E, 0xE7, 0xB7)
RED     = RGBColor(0xDC, 0x26, 0x26)
BLUE    = RGBColor(0x25, 0x63, 0xEB)
BLUE_L  = RGBColor(0x93, 0xC5, 0xFD)
TEAL    = RGBColor(0x08, 0x91, 0xB2)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Core helpers ─────────────────────────────────────────────────────────────
def ns():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=BG):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def box(slide, x, y, w, h, fill=None, line=None, lw=Pt(1.2)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = lw
    else: s.line.fill.background()
    return s

def txt(slide, x, y, w, h, text, sz=Pt(18), col=OFF_W, bold=False,
        italic=False, align=PP_ALIGN.LEFT, font="Calibri Light", wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = sz; r.font.color.rgb = col
    r.font.bold = bold; r.font.italic = italic; r.font.name = font
    return tb

def hline(slide, y, lpad=Inches(0.4), rpad=Inches(0.4), col=BG3):
    box(slide, lpad, y, W - lpad - rpad, Emu(16000), fill=col)

def slide_header(slide, title, sub=None):
    """Standard slide header: amber accent bar + title + optional subtitle + separator"""
    box(slide, Inches(0.35), Inches(0.3), Inches(0.07), Inches(0.7), fill=AMBER)
    txt(slide, Inches(0.55), Inches(0.22), Inches(12.3), Inches(0.72),
        title, sz=Pt(30), col=AMBER, bold=True, font="Calibri")
    if sub:
        txt(slide, Inches(0.55), Inches(0.94), Inches(12.3), Inches(0.3),
            sub, sz=Pt(13), col=MUTED, font="Calibri Light")
    hline(slide, Inches(1.15))

def bullets(slide, items, x, y, w, h, sz=Pt(16), gap=Pt(5)):
    """items: list of str or (str, level) tuples"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        lvl = 0
        if isinstance(item, tuple): text, lvl = item
        else: text = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run()
        if lvl == 0:
            r.text = "▸  " + text; r.font.size = sz; r.font.color.rgb = OFF_W
        else:
            r.text = "   •  " + text; r.font.size = Pt(14); r.font.color.rgb = MUTED
        r.font.name = "Calibri Light"
        p.space_after = gap
    return tb

def checks(slide, items, x, y, w, sz=Pt(16)):
    tb = slide.shapes.add_textbox(x, y, w, Inches(5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = "✓  " + item
        r.font.size = sz; r.font.color.rgb = GREEN_L; r.font.name = "Calibri Light"
        p.space_after = Pt(6)
    return tb

def label_val(slide, lbl, val, lx, vx, y, w, sz=Pt(14)):
    txt(slide, lx, y, Inches(2.5), Inches(0.45), lbl, sz=sz, col=MUTED, font="Calibri Light")
    txt(slide, vx, y, w, Inches(0.45), val, sz=sz, col=WHITE, bold=True, font="Calibri")

def badge(slide, x, y, w, h, text, fc=AMBER, tc=BG, sz=Pt(12), bold=True):
    box(slide, x, y, w, h, fill=fc, line=None)
    txt(slide, x, y, w, h, text, sz=sz, col=tc, bold=bold, align=PP_ALIGN.CENTER, font="Calibri")

def step_box(slide, x, y, w, h, num, title, body=None, fill=BG3, accent=AMBER):
    box(slide, x, y, w, h, fill=fill, line=accent, lw=Emu(9525))
    # Number circle
    box(slide, x + Emu(50000), y + Emu(50000), Inches(0.4), Inches(0.4), fill=accent)
    txt(slide, x + Emu(50000), y + Emu(50000), Inches(0.4), Inches(0.4),
        str(num), sz=Pt(13), col=BG, bold=True, align=PP_ALIGN.CENTER, font="Calibri")
    txt(slide, x + Inches(0.55), y + Emu(55000), w - Inches(0.65), Inches(0.42),
        title, sz=Pt(13), col=AMBER_L, bold=True, font="Calibri")
    if body:
        txt(slide, x + Inches(0.55), y + Inches(0.52), w - Inches(0.65), h - Inches(0.6),
            body, sz=Pt(11), col=MUTED, font="Calibri Light")

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — BÌA
# ═════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s, BG)

box(s, 0, 0, W, Inches(0.07), fill=AMBER)                       # top stripe
box(s, 0, Inches(0.07), W, Inches(0.75), fill=BG3)             # header band
txt(s, Inches(0.5), Inches(0.1), W - Inches(1), Inches(0.38),
    "TRƯỜNG ĐẠI HỌC TRÀ VINH  ·  KHOA KỸ THUẬT VÀ CÔNG NGHỆ",
    sz=Pt(13), col=MUTED, align=PP_ALIGN.CENTER, font="Calibri Light")
txt(s, Inches(0.5), Inches(0.46), W - Inches(1), Inches(0.34),
    "ĐỒ ÁN THỰC TẬP CƠ SỞ NGÀNH  ·  CHUYÊN NGÀNH CNTT",
    sz=Pt(12), col=AMBER, bold=True, align=PP_ALIGN.CENTER, font="Calibri")

txt(s, Inches(0.5), Inches(1.15), W - Inches(1), Inches(0.72),
    "TÌM HIỂU AMAZON ELASTIC COMPUTE CLOUD",
    sz=Pt(36), col=AMBER, bold=True, align=PP_ALIGN.CENTER, font="Calibri")
txt(s, Inches(0.5), Inches(1.9), W - Inches(1), Inches(0.52),
    "VÀ THIẾT KẾ WEBSITE BÁN SẢN PHẨM ĐIỆN MÁY",
    sz=Pt(22), col=OFF_W, bold=False, align=PP_ALIGN.CENTER, font="Calibri Light")

box(s, Inches(4.0), Inches(2.67), Inches(5.33), Emu(16000), fill=AMBER)

badges = ["Next.js 16", "Payload CMS", "MongoDB", "AWS EC2"]
bw = Inches(1.9); bh = Inches(0.42); gap = Emu(120000)
total = len(badges) * bw + (len(badges)-1) * gap
sx = (W - total) / 2
for i, b in enumerate(badges):
    bx = sx + i*(bw + gap)
    box(s, bx, Inches(2.95), bw, bh, fill=BG3, line=AMBER, lw=Emu(9525))
    txt(s, bx, Inches(2.96), bw, bh, b, sz=Pt(13), col=AMBER,
        bold=True, align=PP_ALIGN.CENTER, font="Calibri")

info_box_x = Inches(3.5); info_box_y = Inches(3.8)
box(s, info_box_x, info_box_y, Inches(6.33), Inches(2.65), fill=BG3)
box(s, info_box_x, info_box_y, Inches(0.06), Inches(2.65), fill=AMBER)
rows = [("Sinh viên:", "Nguyễn Chí Tôn"),
        ("MSSV:",      "170123297"),
        ("Lớp:",       "DX23TT10"),
        ("GVHD:",      "ThS. Nguyễn Hoàng Duy Thiện")]
for i, (lbl, val) in enumerate(rows):
    ry = info_box_y + Inches(0.15) + i * Inches(0.6)
    label_val(s, lbl, val, info_box_x + Inches(0.2), info_box_x + Inches(2.0),
              ry, Inches(4.2))

txt(s, 0, Inches(6.85), W, Inches(0.4), "Trà Vinh, tháng 06 năm 2026",
    sz=Pt(11), col=MUTED, align=PP_ALIGN.CENTER, font="Calibri Light")
box(s, 0, Inches(7.3), W, Inches(0.2), fill=AMBER)

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — MỤC LỤC
# ═════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s, BG)
slide_header(s, "Nội dung thuyết trình")

toc = [
    ("01", "Giới thiệu đề tài"),
    ("02", "Cơ sở lý thuyết – Công nghệ sử dụng"),
    ("03", "Phân tích & Thiết kế hệ thống"),
    ("04", "Cài đặt và triển khai trên AWS EC2"),
    ("05", "Kết quả & Demo"),
    ("06", "Kết luận & Hướng phát triển"),
]

cols = 2; rows_per_col = 3
for i, (num, title) in enumerate(toc):
    col = i // rows_per_col
    row = i % rows_per_col
    ix = Inches(0.5) + col * Inches(6.5)
    iy = Inches(1.5) + row * Inches(1.8)
    iw = Inches(6.0); ih = Inches(1.5)
    box(s, ix, iy, iw, ih, fill=BG3, line=AMBER, lw=Emu(6350))
    box(s, ix, iy, Inches(0.8), ih, fill=AMBER)
    txt(s, ix, iy, Inches(0.8), ih, num, sz=Pt(24), col=BG, bold=True,
        align=PP_ALIGN.CENTER, font="Calibri")
    txt(s, ix + Inches(0.95), iy + Inches(0.45), iw - Inches(1.1), Inches(0.65),
        title, sz=Pt(17), col=WHITE, bold=True, font="Calibri")

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — LÝ DO CHỌN ĐỀ TÀI
# ═════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s, BG)
slide_header(s, "Lý do chọn đề tài")

sections = [
    (TEAL,    "Thực trạng",
     ["E-commerce VN tăng trưởng 25%/năm",
      "Điện Máy Xanh, Nguyễn Kim đạt doanh thu nghìn tỷ",
      "Website là kênh bán hàng thiết yếu"]),
    (AMBER,   "Xu hướng",
     ["Cloud Computing là hạ tầng tất yếu",
      "AWS chiếm 31% thị phần cloud toàn cầu",
      "Deploy server không cần mua phần cứng"]),
    (GREEN,   "Mục tiêu đề tài",
     ["Tìm hiểu AWS EC2 – máy chủ đám mây",
      "Xây dựng website điện máy hoàn chỉnh",
      "Triển khai trên môi trường thực tế"]),
]

for i, (accent, heading, points) in enumerate(sections):
    cx = Inches(0.4) + i * Inches(4.35)
    cy = Inches(1.35); cw = Inches(4.1); ch = Inches(5.7)
    box(s, cx, cy, cw, ch, fill=BG3)
    box(s, cx, cy, cw, Inches(0.06), fill=accent)
    txt(s, cx + Inches(0.2), cy + Inches(0.18), cw - Inches(0.4), Inches(0.5),
        heading, sz=Pt(18), col=accent, bold=True, font="Calibri")
    hline(s, cy + Inches(0.62), lpad=cx + Inches(0.15), rpad=W - cx - cw + Inches(0.15),
          col=BG2)
    for j, pt in enumerate(points):
        py = cy + Inches(0.85) + j * Inches(1.45)
        box(s, cx + Inches(0.15), py, Inches(0.4), Inches(0.4), fill=accent)
        txt(s, cx + Inches(0.15), py, Inches(0.4), Inches(0.4),
            str(j+1), sz=Pt(14), col=BG, bold=True, align=PP_ALIGN.CENTER, font="Calibri")
        txt(s, cx + Inches(0.65), py - Emu(15000), cw - Inches(0.85), Inches(0.85),
            pt, sz=Pt(14), col=OFF_W, font="Calibri Light", wrap=True)

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — CÔNG NGHỆ SỬ DỤNG
# ═════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s, BG)
slide_header(s, "Công nghệ sử dụng", "Technology Stack")

techs = [
    (AMBER,   "Next.js 15",       "Framework React fullstack – App Router, SSR, TypeScript"),
    (TEAL,    "Payload CMS 3",    "Headless CMS – Admin panel, REST API, auto collections"),
    (GREEN,   "MongoDB 7",        "NoSQL document database – Atlas hoặc local instance"),
    (RED,     "AWS EC2",          "Virtual machine t2.micro – Ubuntu 24.04 LTS, Free Tier"),
    (BLUE,    "Nginx",            "Reverse proxy – port 80 → localhost:3000"),
    (MUTED,   "PM2",              "Process manager – auto restart, startup on boot"),
    (AMBER_L, "Xshell / Xftp",   "SSH terminal & SFTP file transfer client"),
]

for i, (color, name, desc) in enumerate(techs):
    col = i % 2; row = i // 2
    ix = Inches(0.4) + col * Inches(6.5)
    iy = Inches(1.45) + row * Inches(1.38)
    iw = Inches(6.1); ih = Inches(1.22)
    if i == 6:
        ix = Inches(3.6); iw = Inches(6.1)

    box(s, ix, iy, iw, ih, fill=BG3)
    box(s, ix, iy, Inches(0.06), ih, fill=color)
    box(s, ix + Inches(0.16), iy + Inches(0.12), Inches(0.95), Inches(0.42),
        fill=color)
    txt(s, ix + Inches(0.16), iy + Inches(0.12), Inches(0.95), Inches(0.42),
        name, sz=Pt(11), col=BG, bold=True, align=PP_ALIGN.CENTER, font="Calibri")
    txt(s, ix + Inches(1.25), iy + Inches(0.15), iw - Inches(1.4), Inches(0.38),
        name, sz=Pt(16), col=color, bold=True, font="Calibri")
    txt(s, ix + Inches(1.25), iy + Inches(0.56), iw - Inches(1.4), Inches(0.6),
        desc, sz=Pt(12), col=MUTED, font="Calibri Light")

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — AMAZON EC2
# ═════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s, BG)
slide_header(s, "Amazon EC2 là gì?", "Elastic Compute Cloud – Virtual Machine trên đám mây AWS")

txt(s, Inches(0.4), Inches(1.35), Inches(8.2), Inches(0.45),
    "Máy chủ ảo (Virtual Machine) chạy trên hạ tầng AWS – khởi tạo trong vài phút",
    sz=Pt(15), col=OFF_W, font="Calibri Light")

# Left: advantages
left_items = [
    "Khởi tạo trong vài phút, ngay lập tức",
    "Scale up/down dễ dàng theo nhu cầu",
    "Free Tier: t2.micro miễn phí 12 tháng",
    "Uptime 99.99% – SLA Amazon",
    "Security Group = Firewall ảo linh hoạt",
]
for i, item in enumerate(left_items):
    iy = Inches(2.0) + i * Inches(0.9)
    box(s, Inches(0.4), iy + Emu(80000), Inches(0.35), Inches(0.35), fill=AMBER)
    txt(s, Inches(0.4), iy + Emu(80000), Inches(0.35), Inches(0.35),
        str(i+1), sz=Pt(12), col=BG, bold=True, align=PP_ALIGN.CENTER, font="Calibri")
    txt(s, Inches(0.9), iy, Inches(5.8), Inches(0.9), item,
        sz=Pt(15), col=OFF_W, font="Calibri Light")

# Right: components
box(s, Inches(7.2), Inches(1.9), Inches(5.7), Inches(5.3), fill=BG3)
txt(s, Inches(7.4), Inches(2.0), Inches(5.3), Inches(0.45),
    "Cấu hình sử dụng trong đồ án",
    sz=Pt(14), col=AMBER, bold=True, font="Calibri")

comps = [
    ("Instance",      "t2.micro – 1 vCPU, 1 GB RAM"),
    ("OS / AMI",      "Ubuntu Server 24.04 LTS"),
    ("Security Group","Port 22 (SSH) · 80 (HTTP) · 443 (HTTPS)"),
    ("Key Pair",      "RSA .pem – xác thực SSH không cần password"),
    ("Storage",       "20 GB EBS gp3 SSD"),
    ("Region",        "ap-southeast-1 (Singapore)"),
]
for i, (k, v) in enumerate(comps):
    ry = Inches(2.55) + i * Inches(0.77)
    txt(s, Inches(7.4), ry, Inches(2.0), Inches(0.45),
        k, sz=Pt(13), col=AMBER_L, bold=True, font="Calibri")
    txt(s, Inches(9.55), ry, Inches(3.2), Inches(0.7),
        v, sz=Pt(13), col=OFF_W, font="Calibri Light")

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — KIẾN TRÚC HỆ THỐNG
# ═════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s, BG)
slide_header(s, "Kiến trúc hệ thống", "System Architecture")

# Architecture boxes (left to right flow)
arch = [
    (TEAL,   "BROWSER",       "Port 80 / 443"),
    (AMBER,  "NGINX",         "Reverse Proxy\nport 80 → :3000"),
    (BG3,    "Next.js App",   "PM2 · Port 3000\nApp Router + API"),
    (GREEN,  "MongoDB",       "Port 27017\nLocal instance"),
]
for i, (col, title, body) in enumerate(arch):
    ax = Inches(0.5) + i * Inches(3.1)
    ay = Inches(2.4); aw = Inches(2.75); ah = Inches(2.0)
    box(s, ax, ay, aw, ah, fill=col if col != BG3 else BG3,
        line=col if col == BG3 else None, lw=Emu(12700))
    txt(s, ax, ay, aw, Inches(0.45),
        title, sz=Pt(14), col=BG if col != BG3 else AMBER,
        bold=True, align=PP_ALIGN.CENTER, font="Calibri")
    txt(s, ax + Inches(0.1), ay + Inches(0.5), aw - Inches(0.2), Inches(1.35),
        body, sz=Pt(12), col=BG if col != BG3 else MUTED,
        align=PP_ALIGN.CENTER, font="Calibri Light")
    if i < len(arch)-1:
        ax2 = ax + aw + Emu(20000)
        txt(s, ax2, ay + Inches(0.8), Emu(180000), Inches(0.45),
            "→", sz=Pt(26), col=AMBER, align=PP_ALIGN.CENTER, font="Calibri")

# CMS sub-note
box(s, Inches(6.4), Inches(4.7), Inches(3.2), Inches(1.75), fill=BG3, line=AMBER_D, lw=Emu(6350))
txt(s, Inches(6.55), Inches(4.8), Inches(2.9), Inches(0.38),
    "Payload CMS (trong Next.js)", sz=Pt(13), col=AMBER, bold=True, font="Calibri")
for j, route in enumerate(["/admin → Quản trị", "/api/* → REST API", "/api/orders → Đặt hàng"]):
    txt(s, Inches(6.55), Inches(5.22) + j * Inches(0.37), Inches(2.9), Inches(0.36),
        "▸ " + route, sz=Pt(12), col=MUTED, font="Calibri Light")

# Routes list
box(s, Inches(0.4), Inches(4.7), Inches(5.7), Inches(2.55), fill=BG3)
txt(s, Inches(0.55), Inches(4.82), Inches(5.4), Inches(0.38),
    "Routes frontend", sz=Pt(13), col=AMBER, bold=True, font="Calibri")
routes = ["/ – Trang chủ", "/san-pham – Danh sách SP", "/san-pham/[slug] – Chi tiết",
          "/tim-kiem – Tìm kiếm", "/gio-hang, /dat-hang – Mua hàng", "/lien-he – Liên hệ"]
for j, r in enumerate(routes):
    txt(s, Inches(0.55), Inches(5.25) + j * Inches(0.32), Inches(5.4), Inches(0.32),
        "▸ " + r, sz=Pt(11), col=MUTED, font="Calibri Light")

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — THIẾT KẾ DATABASE
# ═════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s, BG)
slide_header(s, "Thiết kế cơ sở dữ liệu", "MongoDB – 6 Collections")

cols_data = [
    (AMBER,  "products",          ["_id (PK)", "name, slug (unique)", "category → FK categories",
                                   "price, originalPrice", "brand, inStock, featured",
                                   "images[] → FK media", "specifications [{key,val}]"]),
    (TEAL,   "categories",        ["_id (PK)", "name, slug (unique)", "description",
                                   "image → FK media"]),
    (GREEN,  "orders",            ["_id (PK)", "orderCode (TS-XXXX)", "customerName, phone, address",
                                   "items [{id,name,price,qty}]", "totalAmount",
                                   "paymentMethod: cod|transfer", "status: pending|confirmed|done"]),
    (BLUE,   "media",             ["_id (PK)", "filename, url", "mimeType", "width, height, filesize"]),
    (MUTED,  "users",             ["_id (PK)", "email (unique)", "password (hashed)", "role: admin|user"]),
    (RED,    "contact-messages",  ["_id (PK)", "name, email, phone", "subject, message",
                                   "status: new|read|replied"]),
]

for i, (color, cname, fields) in enumerate(cols_data):
    col = i % 3; row = i // 3
    cx = Inches(0.35) + col * Inches(4.35)
    cy = Inches(1.35) + row * Inches(2.85)
    cw = Inches(4.1); ch = Inches(2.68)
    box(s, cx, cy, cw, ch, fill=BG3)
    box(s, cx, cy, cw, Inches(0.44), fill=color)
    txt(s, cx + Inches(0.12), cy, cw - Inches(0.24), Inches(0.44),
        cname, sz=Pt(15), col=BG if color != MUTED else WHITE,
        bold=True, align=PP_ALIGN.CENTER, font="Calibri")
    for j, f in enumerate(fields):
        txt(s, cx + Inches(0.18), cy + Inches(0.52) + j * Inches(0.33),
            cw - Inches(0.3), Inches(0.33),
            "• " + f, sz=Pt(11), col=MUTED, font="Calibri Light")

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — PAYLOAD CMS ADMIN PANEL
# ═════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s, BG)
slide_header(s, "Payload CMS – Admin Panel", "Giao diện quản trị tại /admin")

features = [
    (AMBER,  "Products",         "Thêm/sửa/xóa sản phẩm, upload nhiều ảnh,\nnhập thông số kỹ thuật, đánh dấu nổi bật"),
    (TEAL,   "Categories",       "Quản lý 8 danh mục: Tivi, Điều hòa,\nTủ lạnh, Máy giặt, Laptop..."),
    (RED,    "Orders",           "Xem đơn hàng mới, cập nhật trạng thái\npending → confirmed → done"),
    (BLUE,   "Media",            "Upload & quản lý hình ảnh sản phẩm,\nauto resize và lưu URL"),
    (MUTED,  "Users",            "Tạo tài khoản admin, phân quyền,\nxác thực JWT tích hợp sẵn"),
    (GREEN,  "Contact Messages", "Xem tin nhắn từ trang Liên hệ,\ncập nhật trạng thái new|read|replied"),
]

for i, (color, name, desc) in enumerate(features):
    col = i % 2; row = i // 2
    fx = Inches(0.4) + col * Inches(6.5)
    fy = Inches(1.45) + row * Inches(1.98)
    fw = Inches(6.1); fh = Inches(1.8)
    box(s, fx, fy, fw, fh, fill=BG3)
    box(s, fx, fy, Inches(0.06), fh, fill=color)
    txt(s, fx + Inches(0.2), fy + Inches(0.18), fw - Inches(0.3), Inches(0.45),
        name, sz=Pt(17), col=color, bold=True, font="Calibri")
    txt(s, fx + Inches(0.2), fy + Inches(0.68), fw - Inches(0.3), Inches(1.0),
        desc, sz=Pt(13), col=MUTED, font="Calibri Light")

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — GIAO DIỆN WEBSITE
# ═════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s, BG)
slide_header(s, "Giao diện website", "Demo – TechShop Frontend")

# Left: screenshot placeholder
box(s, Inches(0.4), Inches(1.35), Inches(6.0), Inches(3.0), fill=BG3, line=AMBER, lw=Emu(9525))
txt(s, Inches(0.4), Inches(1.35), Inches(6.0), Inches(3.0),
    "Screenshot\nTrang Chủ",
    sz=Pt(18), col=MUTED, align=PP_ALIGN.CENTER, font="Calibri Light")

box(s, Inches(0.4), Inches(4.5), Inches(6.0), Inches(2.7), fill=BG3, line=TEAL, lw=Emu(9525))
txt(s, Inches(0.4), Inches(4.5), Inches(6.0), Inches(2.7),
    "Screenshot\nTrang Sản Phẩm",
    sz=Pt(18), col=MUTED, align=PP_ALIGN.CENTER, font="Calibri Light")

# Right: descriptions
txt(s, Inches(6.7), Inches(1.35), Inches(6.2), Inches(0.45),
    "Trang Chủ (/)", sz=Pt(18), col=AMBER, bold=True, font="Calibri")
home_pts = ["Hero banner – TechShop điện máy chính hãng",
            "4 cam kết: Giao hàng · Chính hãng · Bảo hành · Đổi trả",
            "Grid 6 danh mục nhanh (Tivi, Điều hòa, Tủ lạnh...)",
            "Grid 8 sản phẩm nổi bật (featured=true)"]
for j, p in enumerate(home_pts):
    txt(s, Inches(6.7), Inches(1.88) + j*Inches(0.42), Inches(6.2), Inches(0.42),
        "▸  " + p, sz=Pt(13), col=OFF_W, font="Calibri Light")

txt(s, Inches(6.7), Inches(4.5), Inches(6.2), Inches(0.45),
    "Trang Sản Phẩm (/san-pham)", sz=Pt(18), col=TEAL, bold=True, font="Calibri")
prod_pts = ["Sidebar lọc: danh mục + giá slider",
            "Grid 4 cột – card sản phẩm với % giảm giá",
            "Filter theo ?danh-muc= query param",
            "Sắp xếp: mới nhất, giá tăng/giảm"]
for j, p in enumerate(prod_pts):
    txt(s, Inches(6.7), Inches(5.05) + j*Inches(0.42), Inches(6.2), Inches(0.42),
        "▸  " + p, sz=Pt(13), col=OFF_W, font="Calibri Light")

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — CHI TIẾT SP & TÌM KIẾM & GIỎ HÀNG
# ═════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s, BG)
slide_header(s, "Chi tiết sản phẩm & Đặt hàng", "Demo – Product Detail & Checkout Flow")

sections_10 = [
    (AMBER, "Chi tiết SP (/san-pham/[slug])",
     ["Gallery ảnh sản phẩm + zoom",
      "Giá đỏ nổi bật, giá gốc gạch ngang, % tiết kiệm",
      "Bảng thông số kỹ thuật (specifications)",
      "Nút thêm giỏ hàng, thêm yêu thích"]),
    (TEAL, "Giỏ hàng & Đặt hàng (/gio-hang, /dat-hang)",
     ["CartContext – lưu state giỏ hàng React",
      "Cập nhật số lượng, xóa sản phẩm",
      "Form checkout: tên, SĐT, địa chỉ, ghi chú",
      "Thanh toán COD hoặc Chuyển khoản",
      "POST /api/orders → MongoDB orders collection"]),
    (GREEN, "Tìm kiếm & Yêu thích",
     ["Tìm kiếm fulltext: /tim-kiem?q=Samsung",
      "Filter kết quả theo danh mục",
      "Yêu thích lưu localStorage (không cần login)",
      "Badge đếm số lượng trên Header icon"]),
]

for i, (color, title, pts) in enumerate(sections_10):
    sx = Inches(0.35) + i * Inches(4.35)
    sy = Inches(1.35)
    sw = Inches(4.1)
    box(s, sx, sy, sw, Inches(5.9), fill=BG3)
    box(s, sx, sy, sw, Inches(0.06), fill=color)
    txt(s, sx + Inches(0.15), sy + Inches(0.15), sw - Inches(0.3), Inches(0.5),
        title, sz=Pt(14), col=color, bold=True, font="Calibri", wrap=True)
    for j, pt in enumerate(pts):
        txt(s, sx + Inches(0.18), sy + Inches(0.82) + j*Inches(0.98),
            sw - Inches(0.3), Inches(0.85),
            "▸  " + pt, sz=Pt(13), col=OFF_W, font="Calibri Light", wrap=True)

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — QUY TRÌNH DEPLOY EC2
# ═════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s, BG)
slide_header(s, "Quy trình triển khai lên AWS EC2", "6 bước deploy production")

steps = [
    ("Tạo EC2 Instance",     "Ubuntu 24.04 · t2.micro · Security Group\n22/80/443 · Key Pair RSA"),
    ("Kết nối SSH (Xshell)", "ubuntu@<Elastic-IP> -i techshop-key.pem\nThêm fingerprint, xác nhận kết nối"),
    ("Cài đặt môi trường",   "Node.js 20 LTS · MongoDB 7 · Nginx · PM2\nTạo 2GB swap cho t2.micro 1GB RAM"),
    ("Upload code (Xftp)",   "SFTP → /home/ubuntu/dien-may-website\nThêm file .env với MONGODB_URI, keys"),
    ("Build & chạy app",     "npm install && npm run build (swap ~3 phút)\npm2 start npm -- start · pm2 startup"),
    ("Nginx Reverse Proxy",  "proxy_pass http://127.0.0.1:3000\nsystemctl enable nginx · test cấu hình"),
]

cols_n = 3
for i, (title, body) in enumerate(steps):
    col = i % cols_n; row = i // cols_n
    step_box(s,
        x = Inches(0.35) + col * Inches(4.35),
        y = Inches(1.35) + row * Inches(2.85),
        w = Inches(4.1), h = Inches(2.7),
        num = i+1, title = title, body = body)

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — SAO LƯU & PHỤC HỒI
# ═════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s, BG)
slide_header(s, "Sao lưu & Phục hồi dữ liệu", "Backup & Restore MongoDB")

# Left: backup steps
steps12 = [
    (AMBER, "1. Backup thủ công (Xshell)",
     "mongodump --db dien-may-website\n  --out /home/ubuntu/backup/$(date +%Y%m%d)"),
    (TEAL, "2. Download file (Xftp)",
     "SFTP → /home/ubuntu/backup/ → kéo\nthư mục backup về máy local"),
    (GREEN, "3. Backup tự động hàng ngày",
     "crontab -e\n0 2 * * *  mongodump --db ... --out ..."),
    (RED, "4. Restore khi sự cố",
     "mongorestore --db dien-may-website\n  /home/ubuntu/backup/20260601/"),
]

for i, (color, title, cmd) in enumerate(steps12):
    sy = Inches(1.4) + i * Inches(1.42)
    box(s, Inches(0.35), sy, Inches(7.0), Inches(1.28), fill=BG3, line=color, lw=Emu(9525))
    box(s, Inches(0.35), sy, Inches(0.06), Inches(1.28), fill=color)
    txt(s, Inches(0.55), sy + Inches(0.1), Inches(6.6), Inches(0.42),
        title, sz=Pt(14), col=color, bold=True, font="Calibri")
    txt(s, Inches(0.55), sy + Inches(0.56), Inches(6.6), Inches(0.65),
        cmd, sz=Pt(12), col=MUTED, font="Courier New")

# Right: summary
box(s, Inches(7.7), Inches(1.4), Inches(5.25), Inches(5.7), fill=BG3)
txt(s, Inches(7.9), Inches(1.55), Inches(4.85), Inches(0.45),
    "Lưu ý khi sao lưu", sz=Pt(16), col=AMBER, bold=True, font="Calibri")
notes = ["Backup trước khi update code hoặc data",
         "File backup nên lưu ít nhất 7 ngày gần nhất",
         "Kiểm tra dung lượng EBS (df -h) định kỳ",
         "Nên thêm S3 sync cho backup offsite",
         "mongodump tạo thư mục BSON có thể restore y hệt"]
for j, n in enumerate(notes):
    txt(s, Inches(7.9), Inches(2.1) + j*Inches(0.95), Inches(4.85), Inches(0.85),
        "▸  " + n, sz=Pt(13), col=OFF_W, font="Calibri Light", wrap=True)

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — KẾT QUẢ ĐẠT ĐƯỢC
# ═════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s, BG)
slide_header(s, "Kết quả đạt được", "Tổng kết những gì đã thực hiện")

results = [
    ("AWS EC2",           "Tạo và quản lý Instance t2.micro thành công",          True),
    ("Website hoạt động", "9 trang frontend + Admin panel + REST API",             True),
    ("Sản phẩm",          "18 sản phẩm, 8 danh mục đã seed dữ liệu thực",         True),
    ("Deploy Production", "Nginx + PM2 chạy ổn định trên Ubuntu 24.04",           True),
    ("Backup / Restore",  "mongodump/mongorestore, crontab tự động",               True),
    ("Xshell / Xftp",     "Kết nối SSH, quản lý file, upload code thành công",     True),
    ("Báo cáo & Slide",   "Tài liệu đầy đủ theo mẫu TVU",                         True),
]

for i, (obj, res, done) in enumerate(results):
    col = i % 2; row = i // 2
    if i == 6: col = 0; row = 3
    rx = Inches(0.35) + col * Inches(6.5)
    ry = Inches(1.4) + row * Inches(1.48)
    rw = Inches(6.1); rh = Inches(1.32)
    box(s, rx, ry, rw, rh, fill=BG3)
    dot_col = GREEN if done else RED
    box(s, rx + Inches(0.15), ry + Inches(0.44), Inches(0.4), Inches(0.4), fill=dot_col)
    txt(s, rx + Inches(0.15), ry + Inches(0.44), Inches(0.4), Inches(0.4),
        "✓" if done else "✗", sz=Pt(14), col=BG, bold=True,
        align=PP_ALIGN.CENTER, font="Calibri")
    txt(s, rx + Inches(0.7), ry + Inches(0.1), rw - Inches(0.85), Inches(0.46),
        obj, sz=Pt(15), col=AMBER_L, bold=True, font="Calibri")
    txt(s, rx + Inches(0.7), ry + Inches(0.58), rw - Inches(0.85), Inches(0.65),
        res, sz=Pt(13), col=OFF_W, font="Calibri Light", wrap=True)

# URL info
box(s, Inches(6.75), Inches(1.4), Inches(6.1), Inches(1.32), fill=BG3, line=AMBER, lw=Emu(9525))
txt(s, Inches(7.0), Inches(1.52), Inches(5.7), Inches(0.45),
    "URL triển khai", sz=Pt(14), col=AMBER, bold=True, font="Calibri")
txt(s, Inches(7.0), Inches(2.0), Inches(5.7), Inches(0.65),
    "http://<Elastic-IP>        →  Website\nhttp://<Elastic-IP>/admin  →  Admin Panel",
    sz=Pt(13), col=MUTED, font="Courier New")

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — KẾT LUẬN & HƯỚNG PHÁT TRIỂN
# ═════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s, BG)
slide_header(s, "Kết luận & Hướng phát triển")

# Left: conclusion
box(s, Inches(0.35), Inches(1.35), Inches(6.0), Inches(5.9), fill=BG3)
box(s, Inches(0.35), Inches(1.35), Inches(6.0), Inches(0.06), fill=AMBER)
txt(s, Inches(0.55), Inches(1.5), Inches(5.6), Inches(0.45),
    "Kết luận", sz=Pt(20), col=AMBER, bold=True, font="Calibri")
concs = [
    "Nắm vững AWS EC2: tạo instance, SSH, deploy",
    "Xây dựng website thương mại điện tử hoàn chỉnh",
    "Áp dụng stack hiện đại: Next.js + Payload CMS + MongoDB",
    "Triển khai Production: Nginx reverse proxy + PM2",
    "Thực hành Backup/Restore dữ liệu thực tế",
    "Sử dụng Xshell/Xftp để quản lý server từ xa",
]
for j, c in enumerate(concs):
    txt(s, Inches(0.55), Inches(2.1) + j*Inches(0.85), Inches(5.6), Inches(0.8),
        "✓  " + c, sz=Pt(14), col=GREEN_L, font="Calibri Light", wrap=True)

# Right: future
box(s, Inches(6.65), Inches(1.35), Inches(6.3), Inches(5.9), fill=BG3)
box(s, Inches(6.65), Inches(1.35), Inches(6.3), Inches(0.06), fill=BLUE)
txt(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(0.45),
    "Hướng phát triển", sz=Pt(20), col=BLUE_L, bold=True, font="Calibri")
futures = [
    (AMBER,  "Thanh toán online",  "Tích hợp VNPay, MoMo, ZaloPay"),
    (TEAL,   "Tài khoản người dùng","Đăng ký/đăng nhập, lịch sử đơn hàng"),
    (GREEN,  "SSL / HTTPS",        "Let's Encrypt – cerbot auto-renew"),
    (BLUE,   "CDN & CloudFront",   "Tăng tốc độ tải ảnh toàn cầu"),
    (RED,    "Monitoring",         "AWS CloudWatch – cảnh báo CPU/RAM"),
]
for j, (color, title, desc) in enumerate(futures):
    fy = Inches(2.1) + j*Inches(0.98)
    box(s, Inches(6.85), fy + Emu(60000), Inches(0.35), Inches(0.35), fill=color)
    txt(s, Inches(7.35), fy + Emu(20000), Inches(5.4), Inches(0.42),
        title, sz=Pt(15), col=color, bold=True, font="Calibri")
    txt(s, Inches(7.35), fy + Inches(0.48), Inches(5.4), Inches(0.4),
        desc, sz=Pt(13), col=MUTED, font="Calibri Light")

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — CẢM ƠN / Q&A
# ═════════════════════════════════════════════════════════════════════════════
s = ns(); bg(s, BG)

box(s, 0, 0, W, Inches(0.07), fill=AMBER)
box(s, 0, Inches(7.43), W, Inches(0.07), fill=AMBER)

txt(s, 0, Inches(1.3), W, Inches(1.0),
    "XIN CẢM ƠN THẦY/CÔ VÀ CÁC BẠN",
    sz=Pt(40), col=AMBER, bold=True, align=PP_ALIGN.CENTER, font="Calibri")
txt(s, 0, Inches(2.35), W, Inches(0.7),
    "ĐÃ LẮNG NGHE!",
    sz=Pt(34), col=WHITE, bold=False, align=PP_ALIGN.CENTER, font="Calibri Light")

box(s, Inches(4.0), Inches(3.25), Inches(5.33), Emu(16000), fill=AMBER)

box(s, Inches(3.5), Inches(3.5), Inches(6.33), Inches(1.3), fill=BG3)
txt(s, Inches(3.5), Inches(3.55), Inches(6.33), Inches(1.2),
    "❓  CÂU HỎI VÀ THẢO LUẬN",
    sz=Pt(24), col=AMBER_L, bold=True, align=PP_ALIGN.CENTER, font="Calibri")

info_14 = [("Sinh viên", "Nguyễn Chí Tôn  –  MSSV: 170123297  –  Lớp: DX23TT10"),
           ("GVHD",      "ThS. Nguyễn Hoàng Duy Thiện")]
for j, (lbl, val) in enumerate(info_14):
    txt(s, 0, Inches(5.3) + j*Inches(0.55), W, Inches(0.5),
        lbl + ":  " + val,
        sz=Pt(14), col=MUTED, align=PP_ALIGN.CENTER, font="Calibri Light")

txt(s, 0, Inches(6.62), W, Inches(0.5),
    "Trà Vinh, tháng 06 năm 2026",
    sz=Pt(12), col=BG3, align=PP_ALIGN.CENTER, font="Calibri Light")

# ── Save ─────────────────────────────────────────────────────────────────────
out = "/Users/tcn9911/Ton/do-an-tvu-2026-06/SLIDE_TechShop_NguyenChiTon.pptx"
prs.save(out)
print(f"Saved → {out}  ({prs.slides.__len__()} slides)")
